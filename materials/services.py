# materials/services.py
import logging
import os
import PyPDF2
import docx
from pptx import Presentation
from core.exceptions import FileProcessingError
from .models import ExamDocument

logger = logging.getLogger(__name__)

class MaterialService:
    """
    A service class for handling business logic related to study materials,
    such as file parsing and text extraction.
    """
    @staticmethod
    def extract_text_from_file(file):
        """
        Extracts text content from a supported file type.
        
        Args:
            file: The Django InMemoryUploadedFile object.
            
        Returns:
            str: The extracted text.
            
        Raises:
            FileProcessingError: If the file is unsupported or processing fails.
        """
        filename = file.name.lower()
        file_ext = os.path.splitext(filename)[1]
        
        try:
            if file_ext == '.pdf':
                pdf_reader = PyPDF2.PdfReader(file)
                text = ''.join(page.extract_text() or '' for page in pdf_reader.pages)
            elif file_ext == '.docx':
                doc = docx.Document(file)
                text = '\n'.join([para.text for para in doc.paragraphs])
            elif file_ext == '.pptx':
                prs = Presentation(file)
                text = ''.join(shape.text + '\n' for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
            elif file_ext == '.txt':
                text = file.read().decode('utf-8', errors='ignore')
            else:
                raise FileProcessingError('Unsupported file type.')
            
            if len(text) > 20000:
                text = text[:20000] + '\n... [truncated]'
                
            return text.strip()
            
        except Exception as e:
            logger.error(f"Error processing file {file.name}: {e}")
            raise FileProcessingError(f"Failed to process file: {e}")