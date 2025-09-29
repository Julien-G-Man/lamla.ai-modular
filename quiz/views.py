from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from django.http import JsonResponse, HttpResponse, FileResponse, HttpResponseRedirect 
from django.views.decorators.csrf import csrf_exempt
from django.views.decorators.http import require_http_methods
from django.utils.decorators import method_decorator
from allauth.account.views import LoginView as AllauthLoginView
from allauth.account.adapter import DefaultAccountAdapter
import PyPDF2
import docx
from urllib.parse import quote
from django.utils import timezone
import base64, pytesseract
from .models import (
    Question, QuestionCache, QuizSession
)
import json
import logging
import os
from django.core.exceptions import ValidationError as DjangoValidationError
from django.template.loader import render_to_string
import uuid
from django.urls import reverse
import uuid

from .services import QuizService
from .question_generator import generate_questions_from_text
from .flashcard_generator import generate_flashcards_from_text
from .quiz_download_utils import handle_quiz_download
from .exam_analyzer import perform_exam_analysis
from core.cookies import set_quiz_preference_cookie, get_quiz_preference_cookie 

logger = logging.getLogger(__name__)

# Standard subjects for validation
STANDARD_SUBJECTS = [
    "Mathematics", "Computer Science", "Engineering", "Biology", 
    "Chemistry", "Physics", "English", "History", "Geography", "Economics"
]

def custom_quiz(request):
    """
    Renders the page for creating a custom quiz.
    """
    # Clear any existing quiz data from session
    if 'quiz_questions' in request.session:
        del request.session['quiz_questions']
    if 'quiz_results' in request.session:
        del request.session['quiz_results']
    
    # READ PREFERENCE COOKIE: Get default number of MCQ questions
    # Get the value, defaulting to 5 if not found
    default_num_mcq = int(get_quiz_preference_cookie(request, 'pref_num_mcq', 5))
    
    context = {
        'default_num_mcq': default_num_mcq
        # You can add other defaults here (e.g., pref_difficulty)
    }    
    
    return render(request, 'quiz/custom_quiz.html', context)


@require_http_methods(["POST"])
def ajax_extract_text(request):
    """Extracts text from uploaded files."""
    if 'slide_file' not in request.FILES:
        return JsonResponse({'error': 'No file uploaded'}, status=400)

    file = request.FILES['slide_file']
    filename = file.name.lower()
    file_ext = os.path.splitext(filename)[1]
    max_size = 10 * 1024 * 1024  # 10MB
    
    if file.size > max_size:
        return JsonResponse({'error': 'File too large (max 10MB)'}, status=400)

    try:
        text = ""
        if file_ext == '.pdf':
            # PDF extraction
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() or ''
                
        elif file_ext == '.docx':
            # DOCX extraction
            doc = docx.Document(file)
            text = '\n'.join([para.text for para in doc.paragraphs])
            
        elif file_ext == '.pptx':
            # PPTX extraction
            try:
                from pptx import Presentation
                prs = Presentation(file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            text += shape.text + '\n'
            except ImportError:
                return JsonResponse({'error': 'PPTX processing not available'}, status=400)
                
        elif file_ext == '.txt':
            # TXT extraction
            text = file.read().decode('utf-8', errors='ignore')
            
        else:
            return JsonResponse({'error': 'Unsupported file type. Please upload PDF, DOCX, PPTX, or TXT.'}, status=400)

        # Clean and validate extracted text
        text = text.strip()
        if not text:
            return JsonResponse({'error': 'No text could be extracted from the file.'}, status=400)
        
        # Limit text length for performance
        if len(text) > 50000:
            text = text[:50000] + '\n... [text truncated due to length]'
            
        return JsonResponse({'text': text})
        
    except Exception as e:
        logger.error(f"Text extraction error: {str(e)}", exc_info=True)
        return JsonResponse({'error': f'Failed to extract text: {str(e)}'}, status=500)


@require_http_methods(["POST"])
def generate_questions(request):
    """
    Generates questions from a user-provided text.
    """
    try:
        # Get form data with proper validation
        study_text = request.POST.get('extractedText', '').strip()
        num_mcq = int(request.POST.get('num_mcq', 5))
        
        # Handle num_short with validation
        num_short_raw = request.POST.get('num_short', '0')
        try:
            num_short = int(num_short_raw)
            if num_short < 0:
                num_short = 0
        except (ValueError, TypeError):
            num_short = 0
        
        # Handle subject selection properly
        subject_select = request.POST.get('subject_select', '')
        custom_subject = request.POST.get('subject', '').strip()
        
        # Determine the actual subject to use
        if subject_select == 'Other' and custom_subject:
            subject = custom_subject
        elif subject_select in STANDARD_SUBJECTS: # Assumes STANDARD_SUBJECTS is defined
            subject = subject_select
        else:
            subject = 'General'  # Default subject
        
        difficulty = request.POST.get('difficulty', 'any')
        quiz_time = int(request.POST.get('quiz_time', 10))
        
        # Store uploaded file name if available
        uploaded_file_name = ''
        if 'slide_file' in request.FILES:
            uploaded_file_name = request.FILES['slide_file'].name
        else:
            # Check if we have a file name from previous upload
            uploaded_file_name = request.POST.get('uploaded_file_name', '')
        
        logger.info(f"Quiz generation request - Text length: {len(study_text)}, MCQ: {num_mcq}, Short: {num_short}, Subject: {subject}")
        
        # Validate inputs
        error_message = None
        if not study_text or len(study_text) < 30:
            error_message = 'Please provide at least 30 characters of study material.'
            logger.warning(f"Quiz generation failed - insufficient text length: {len(study_text)}")
        elif num_mcq <= 0 and num_short <= 0:
            error_message = 'Please select at least one question type (MCQ or Short Answer).'
        elif num_mcq > 20 or num_short > 10:
            error_message = 'Maximum questions exceeded: MCQ (20 max), Short Answer (10 max).'
        elif quiz_time < 1 or quiz_time > 120:
            error_message = 'Quiz time must be between 1 and 120 minutes.'
        
        if error_message:
            messages.error(request, error_message)
            return render(request, 'quiz/custom_quiz.html', {
                'subject': subject,
                'subject_select': subject_select,
                'num_mcq': num_mcq,
                'num_short': num_short,
                'difficulty': difficulty,
                'study_text': study_text,
                'uploaded_file_name': uploaded_file_name,
                'quiz_time': quiz_time,
            })
        
        # Generate questions
        quiz_results = None
        try:
            logger.info("Calling generate_questions_from_text...")
            # Assumes QuizService.generate_quiz is a synchronous blocking call
            quiz_results = QuizService.generate_quiz(
                study_text, 
                num_mcq, 
                num_short, 
                subject, 
                difficulty
            )
            logger.info(f"Quiz generation completed - Results keys: {quiz_results.keys() if quiz_results else 'None'}")
            
            if not quiz_results:
                error_message = 'Quiz generation returned no results. Please try again.'
                logger.error("Quiz generation returned None")
            else:
                # Validate and sanitize the results
                mcq_questions = quiz_results.get('mcq_questions', [])
                short_questions = quiz_results.get('short_questions', [])
                
                # Ensure questions have required fields
                for i, q in enumerate(mcq_questions):
                    if not isinstance(q, dict):
                        mcq_questions[i] = {'question': str(q), 'options': [], 'answer': ''}
                    else:
                        q['question'] = q.get('question', '')
                        q['options'] = q.get('options', [])
                        q['answer'] = q.get('answer', '')
                
                for i, q in enumerate(short_questions):
                    if not isinstance(q, dict):
                        short_questions[i] = {'question': str(q), 'answer': ''}
                    else:
                        q['question'] = q.get('question', '')
                        q['answer'] = q.get('answer', '')
                
                # 1. Strictly limit the lists to the requested counts.
                mcq_questions = mcq_questions[:num_mcq]
                short_questions = short_questions[:num_short]
                
                # 2. Update the quiz_results dictionary with the now-truncated lists.
                quiz_results['mcq_questions'] = mcq_questions
                quiz_results['short_questions'] = short_questions
                quiz_results['subject'] = subject
                
                # Log actual counts (now using the final, strictly sliced lists)
                actual_mcq = len(mcq_questions)
                actual_short = len(short_questions)
                logger.info(f"Generated: {actual_mcq} MCQ, {actual_short} Short Answer questions")
                
                if actual_mcq == 0 and actual_short == 0:
                    error_message = 'No questions could be generated. Please try with different or more detailed content.'
                    logger.error("No questions generated after processing")

        except Exception as e:
            logger.error(f"Quiz generation error: {str(e)}", exc_info=True)
            error_message = f"Quiz generation failed: {str(e)}"
        
        if error_message or not quiz_results:
            messages.error(request, error_message or 'Unknown error occurred during question generation.')
            return render(request, 'quiz/custom_quiz.html', {
                'subject': subject,
                'subject_select': subject_select,
                'num_mcq': num_mcq,
                'num_short': num_short,
                'difficulty': difficulty,
                'study_text': study_text,
                'uploaded_file_name': uploaded_file_name,
                'quiz_time': quiz_time,
            })
        
        # Store questions in session (now using the strictly sliced local variables)
        request.session['quiz_questions'] = {
            'mcq_questions': mcq_questions,
            'short_questions': short_questions,
        }
        request.session['quiz_time'] = quiz_time
        request.session['uploaded_file_name'] = uploaded_file_name
        request.session['quiz_subject'] = subject
        request.session.modified = True
        
        # Generate a unique quiz ID and store in session
        request.session['quiz_id'] = str(uuid.uuid4())
        
        request.session.modified = True
        
        # Clear previous results
        if 'quiz_results' in request.session:
            del request.session['quiz_results']
        if 'quiz_user_answers' in request.session:
            del request.session['quiz_user_answers']
        
        logger.info("Quiz generation successful - redirecting to quiz page")
        # prepare the redirect response
        response = redirect('quiz')
        
        # SET PREFERENCE COOKIE: Set the user's chosen num_mcq as a default preference
        # This will be read next time they load the custom_quiz page.
        response = set_quiz_preference_cookie(response, 'pref_num_mcq', str(num_mcq))
        return response
        
    except Exception as e:
        logger.error(f"Unexpected error in generate_questions: {str(e)}", exc_info=True)
        messages.error(request, f"An unexpected error occurred: {str(e)}")
        return render(request, 'quiz/custom_quiz.html', {
            'study_text': request.POST.get('extractedText', ''),
            'num_mcq': int(request.POST.get('num_mcq', 5)),
            'num_short': int(request.POST.get('num_short', 0)),
            'difficulty': request.POST.get('difficulty', 'any'),
        })

def quiz(request):
    """
    Renders the quiz page with questions from the session.
    """
    quiz_questions = request.session.get('quiz_questions', {})
    print("SESSION QUESTIONS:", quiz_questions)
    
    # Check if results already exist in the session (meaning the quiz was submitted)
    if 'quiz_results' in request.session and request.session.get('quiz_results'):
         # If results are found, immediately redirect to the results page to show them.
        messages.info(request, "You were redirected to your saved results!")
        return redirect('quiz_results')
    
    quiz_time = request.session.get('quiz_time', 10)
    
    if not quiz_questions:
        messages.error(request, 'No quiz has been generated. Please create a quiz first.')
        logger.info(f"Quiz questions in session: {quiz_questions}")
        return redirect('custom_quiz')
    
    # Validate quiz questions structure
    mcq_questions = quiz_questions.get('mcq_questions', [])
    short_questions = quiz_questions.get('short_questions', [])
    
    if not mcq_questions and not short_questions:
        messages.error(request, 'No questions found in the quiz. Please generate a new quiz.')
        return redirect('custom_quiz')
    
    # Ensure questions have proper structure
    for i, q in enumerate(mcq_questions):
        if not isinstance(q, dict):
            mcq_questions[i] = {'question': f'Question {i+1}', 'options': ['A', 'B', 'C', 'D'], 'answer': 'A'}
        else:
            q.setdefault('question', f'MCQ Question {i+1}')
            q.setdefault('options', ['Option A', 'Option B', 'Option C', 'Option D'])
            q.setdefault('answer', 'A')
    
    for i, q in enumerate(short_questions):
        if not isinstance(q, dict):
            short_questions[i] = {'question': f'Short Answer Question {i+1}', 'answer': 'Sample answer'}
        else:
            q.setdefault('question', f'Short Answer Question {i+1}')
            q.setdefault('answer', 'Sample answer')
    
    # Retrieve the unique qiuz id
    quiz_id = request.session.get('quiz_id')
    
    context = {
        'questions': {
            'mcq_questions': mcq_questions,
            'short_questions': short_questions,
            'subject': request.session.get('quiz_subject', 'General')
        },
        'quiz_time': quiz_time,
        'quiz_id': quiz_id,
        'total_questions': len(mcq_questions) + len(short_questions),
        'mcq_count': len(mcq_questions),
        'short_count': len(short_questions),
    }
    
    return render(request, 'quiz/quiz.html', context)


def quiz_results(request):
    """
    Handles quiz submission, feedback submission, and displays results.
    """
    if request.method == 'POST':
        # --- 1. Handle Feedback Submission (Form data via AJAX) ---
        # This checks for the AJAX POST triggered by the star rating system.
        # It relies on the form data being accessible via request.POST
        if 'rating' in request.POST and 'quiz_id' in request.POST:
            try:
                rating = request.POST.get('rating')
                quiz_id = request.POST.get('quiz_id')
                
                # --- TO DO: Save Feedback to Database ---
                # Example:
                # if request.user.is_authenticated:
                #     QuizFeedback.objects.create(
                #         user=request.user, 
                #         quiz_session_id=quiz_id, 
                #         rating=int(rating)
                #     )
                
                logger.info(f"Feedback received: Quiz ID {quiz_id}, Rating {rating}")
                return JsonResponse({'status': 'success', 'message': 'Feedback received.'}, status=200)

            except Exception as e:
                logger.error(f"Failed to process feedback: {str(e)}")
                return JsonResponse({'status': 'error', 'message': 'Failed to process feedback.'}, status=500)

        # --- 2. Handle Main Quiz Submission (Raw JSON via AJAX) ---
        try:
            # This is the expected place for the main quiz submission payload
            data = json.loads(request.body.decode('utf-8'))
            user_answers = data.get('user_answers', {})
            quiz_questions = request.session.get('quiz_questions', {})
            
            if not quiz_questions:
                return JsonResponse({'status': 'error', 'message': 'No quiz data found. Please generate a quiz first.'}, status=400)
            
            mcq_questions = quiz_questions.get('mcq_questions', [])
            short_questions = quiz_questions.get('short_questions', [])
            
            results = {
                'total': len(mcq_questions) + len(short_questions),
                'correct': 0,
                'details': [],
                'quiz_id': request.session.get('current_quiz_id') # Ensure you have a quiz_id in session
            }

            # Grade MCQ questions
            for idx, q in enumerate(mcq_questions):
                user_ans = user_answers.get(str(idx), '').strip().upper()
                correct_ans = q.get('answer', '').strip()
                
                 # We compare the user's selected letter/value (user_ans) to the correct one (correct_ans.upper())
                is_correct = user_ans == correct_ans
                
                # Get the list of options
                question_options = q.get('options', [])
                
                results['details'].append({
                    'question': q.get('question', ''),
                    'user_answer': user_ans,
                    'correct_answer': correct_ans,
                    'is_correct': is_correct,
                    'type': 'mcq',
                    'options': question_options,
                })
                if is_correct:
                    results['correct'] += 1

            # Grade short answer questions
            for idx, q in enumerate(short_questions):
                # The index for short answers continues after MCQ indices
                short_idx = len(mcq_questions) + idx
                user_ans = user_answers.get(str(short_idx), '').strip()
                expected_ans = q.get('answer', '').strip()
                
                # Use QuizService for grading if available, otherwise simple comparison
                try:
                    # Assuming QuizService is available and correctly imported
                    is_correct = QuizService.grade_short_answer(q.get('question', ''), expected_ans, user_ans)
                except NameError:
                    # Fallback if QuizService is not imported or available
                    is_correct = user_ans.lower() == expected_ans.lower()
                except Exception:
                    # Fallback to simple comparison for any other grading error
                    is_correct = user_ans.lower() == expected_ans.lower()
                
                results['details'].append({
                    'question': q.get('question', ''),
                    'user_answer': user_ans,
                    'correct_answer': expected_ans,
                    'is_correct': is_correct,
                    'type': 'short',
                    'options': []
                })
                if is_correct:
                    results['correct'] += 1
            
            # Calculate percentage
            results['percentage'] = (results['correct'] / results['total'] * 100) if results['total'] > 0 else 0
            
            # Store results in session
            request.session['quiz_results'] = results
            request.session['quiz_user_answers'] = user_answers
            
            # Save to database if user is authenticated
            if request.user.is_authenticated:
                try:
                    quiz_subject = request.session.get('quiz_subject', '')
                    QuizSession.objects.create(
                        user=request.user,
                        subject=quiz_subject,
                        total_questions=results['total'],
                        correct_answers=results['correct'],
                        score_percentage=results['percentage'],
                        questions_data=quiz_questions,
                        user_answers=user_answers
                    )
                    logger.info(f"Quiz results saved for user {request.user.username}")
                except Exception as e:
                    logger.error(f"Failed to save quiz session: {str(e)}")
            
            # Return JSON response with the redirect url
            return JsonResponse({
                'status': 'ok',
                'redirect_url': reverse('quiz_results')
            })
            
        except json.JSONDecodeError:
            # This handles the case where the POST is neither feedback nor valid JSON, 
            # which likely shouldn't happen unless the client side is broken, but it catches the original error.
            logger.error(f"Quiz results error: JSON decode failed. Request body: {request.body.decode('utf-8')}", exc_info=True)
            return JsonResponse({'status': 'error', 'message': 'Invalid data format submitted.'}, status=400)
        
        except Exception as e:
            logger.error(f"Quiz results error: {str(e)}", exc_info=True)
            return JsonResponse({'status': 'error', 'message': str(e)}, status=500)
    
    # --- 3. Handle GET request - display results ---
    else:
        results = request.session.get('quiz_results')
        quiz_questions = request.session.get('quiz_questions', {})
        user_answers = request.session.get('quiz_user_answers', {})
        uploaded_file_name = request.session.get('uploaded_file_name', '')
        
        if not results:
            # If no results in session, redirect to quiz generation page
            # Note: Using messages.error requires the messages middleware to be active
            messages.error(request, 'No quiz results found. Please complete a quiz first.')
            return redirect('custom_quiz')
        
        # Prepare context for results page
        context = {
            # Use 'score' and 'total' from the results dict for presentation
            'score': results.get('correct', 0),
            'total': results.get('total', 0),
            'score_percent': results.get('percentage', 0),
            # Use the full quiz_questions for the detailed review loop
            'results': results, # The results dict contains 'details' which is used in the template
            'uploaded_file_name': uploaded_file_name,
            'subject': request.session.get('quiz_subject', 'General'),
        }
        return render(request, 'quiz/quiz_results.html', context)

def download_quiz_text(request):
    """
    Thin wrapper view that delegates all the heavy lifting to the utility file.
    """
    return handle_quiz_download(request)

def flashcards(request):
    """Renders the flashcards page."""
    return render(request, 'quiz/flashcards.html')


def generate_flashcards(request):
    """Handles flashcard generation."""
    return render(request, 'quiz/flashcards.html')


def test_flashcard_generator(request):
    """Renders the test page for flashcards."""
    return render(request, 'quiz/test_flashcard_generator.html')

def exam_analyzer(request):
    ''' Delegates all work to exam_analyzer.py '''
    return perform_exam_analysis(request, 'quiz/exam_analyzer.html')

