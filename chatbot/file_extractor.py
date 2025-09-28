import os
import logging
import PyPDF2
import docx
from pptx import Presentation

logger = logging.getLogger(__name__)

# Custom exception for clean error handling in the view
class FileExtractionError(Exception):
    """Custom exception for file extraction related errors."""
    pass

MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB
MAX_TEXT_LENGTH = 50000          # 50,000 characters

def extract_text_from_file(file):
    """
    Extracts text content from an uploaded file (PDF, DOCX, PPTX, TXT).
    Performs file size validation and handles format-specific extraction.

    Args:
        file: Django UploadedFile object.

    Returns:
        str: The cleaned and truncated extracted text.

    Raises:
        FileExtractionError: If file validation or extraction fails.
    """
    filename = file.name.lower()
    file_ext = os.path.splitext(filename)[1]
    
    # 1. Size Validation
    if file.size > MAX_FILE_SIZE:
        raise FileExtractionError('File too large (max 10MB)')

    text = ""
    
    try:
        if file_ext == '.pdf':
            if PyPDF2 is None:
                raise FileExtractionError('PDF support missing. Install PyPDF2.')
            
            # PDF extraction using PyPDF2
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() or ''
                
        elif file_ext == '.docx':
            if docx is None:
                raise FileExtractionError('DOCX support missing. Install python-docx.')
            
            # DOCX extraction using python-docx
            doc = docx.Document(file)
            text = '\n'.join([para.text for para in doc.paragraphs])
            
        elif file_ext == '.pptx':
            if Presentation is None:
                raise FileExtractionError('PPTX support missing. Install python-pptx.')
                
            # PPTX extraction using python-pptx
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + '\n'
                        
        elif file_ext == '.txt':
            # TXT extraction
            text = file.read().decode('utf-8', errors='ignore')
            
        else:
            raise FileExtractionError('Unsupported file type. Please upload PDF, DOCX, PPTX, or TXT.')

    except FileExtractionError:
        # Re-raise custom errors directly
        raise
    except Exception as e:
        logger.error(f"Text extraction failed for {filename}: {e}", exc_info=True)
        raise FileExtractionError(f'Failed to extract text: {str(e)}')

    # 2. Clean and Validate Extracted Text
    text = text.strip()
    if not text:
        raise FileExtractionError('No readable text could be extracted from the file.')
    
    # 3. Limit text length
    if len(text) > MAX_TEXT_LENGTH:
        text = text[:MAX_TEXT_LENGTH] + '\n... [text truncated due to length]'
        
    return text
